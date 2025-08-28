import streamlit as st
import os
import google.generativeai as genai
from dotenv import load_dotenv
from PIL import Image
import docx
import pptx
import openpyxl
import io
import requests
import base64

# --- Configuration ---
load_dotenv()
try:
    # It's safer to check for the key and handle its absence
    api_key = os.getenv('GOOGLE_API_KEY')
    if not api_key:
        st.error("GOOGLE_API_KEY not found. Please set it in your .env file.")
        st.stop()
    genai.configure(api_key=api_key)
    # Model for text-based tasks
    text_model = genai.GenerativeModel('gemini-1.5-flash')
except Exception as e:
    st.error(f"Error configuring the API: {e}")
    st.stop()

# --- Helper Functions ---
def extract_text_from_file(uploaded_file):
    """Extracts text from various file types."""
    try:
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        if file_extension == ".docx":
            return "\n".join([para.text for para in docx.Document(uploaded_file).paragraphs])
        elif file_extension == ".pptx":
            prs = pptx.Presentation(uploaded_file)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif file_extension == ".xlsx":
            workbook = openpyxl.load_workbook(uploaded_file)
            return "\n".join([" ".join(str(cell) if cell is not None else "" for cell in row) for sheet in workbook.worksheets for row in sheet.iter_rows(values_only=True)])
        elif file_extension == ".sql":
            return uploaded_file.getvalue().decode("utf-8")
        return "Unsupported file type for text extraction."
    except Exception as e:
        return f"Error extracting text: {e}"

def generate_image(prompt):
    """Generates an image using Google's Imagen model via free Gemini API."""
    try:
        # Use Imagen 3 if available, fallback to Gemini Flash
        image_model = genai.GenerativeModel("imagen-3.0")  # free model via Gemini
        response = image_model.generate_content(
            prompt,
            generation_config={"sample_count": 1},
            stream=False
        )

        # Check response for image
        if response and response.candidates:
            for candidate in response.candidates:
                if candidate.content and candidate.content.parts:
                    for part in candidate.content.parts:
                        if part.mime_type.startswith("image/"):
                            image_data = part.data
                            return Image.open(io.BytesIO(image_data))
        return "Failed to generate image. No valid image returned."

    except Exception as e:
        return f"Error generating image: {e}"


# --- System Prompts for Different Tools ---
SYSTEM_PROMPTS = {
    "General Assistant": "You are an expert helpful digital marketing assistant who loves to explain in detail.",
    "Image Generator": "You are an AI image generation assistant. The user will provide a prompt describing an image they want to create.",
    "Ad Copy Generator": "You are an expert copywriter. Your task is to create compelling ad copy based on the user's request. Focus on headlines, body text, and calls-to-action.",
    "Social Media Post Generator": "You are a social media manager. Create engaging posts for the specified platform, including relevant hashtags and a suitable tone.",
    "Email Campaign Writer": "You are an email marketing specialist. Write effective marketing emails with strong subject lines and clear calls-to-action.",
    "Blog Generator": "You are a professional blog writer. Your task is to generate a well-structured, SEO-friendly blog post based on the user's topic and keywords.",
    "SEO Analyst": "You are an SEO expert. Generate relevant short-tail and long-tail keywords, analyze competitor strategies, and provide on-page SEO suggestions.",
    "Content Improver": "You are an expert content editor. Rewrite and improve the user's text based on their stated goal (e.g., make it more persuasive, simplify it).",
    "AI to Human Text Converter": "You are a skilled novel writer. Your task is to rewrite AI-generated text to sound more natural, engaging, and human-like. Focus on varying sentence structure, using more natural language, and adding a human touch.",
    "Digital Marketing Analyst": "You are a digital marketing analyst. Your role is to analyze data, summarize reports, and provide actionable insights.", 
    'AI Image Generator': 'You are an expert image creator and a social media expert, and know everything about instagram posts, facebook posts, stories, social media marketing, youtube marketing, and other social media platforms where pictures are uploaded. Your task is to generate images as per the user's request'
}

# --- Streamlit App ---
st.set_page_config(page_title="Marketing AI Chat", page_icon="🚀", layout="wide")

st.title("🚀 AI Digital Marketing Assistant")

# --- App Capabilities Summary ---
with st.expander("See what this assistant can do for you"):
    st.markdown("""
    This AI-powered assistant is designed to help you with a wide range of digital marketing tasks. Select a tool from the sidebar to set the AI's persona and get started!

    **Available Tools:**
    - **General Assistant**: Get detailed explanations on any marketing topic.
    - **Image Generator**: Create unique images from a text description.
    - **Ad Copy Generator**: Create compelling ad copy for various platforms.
    - **Social Media Post Generator**: Craft engaging posts tailored for different social channels.
    - **Email Campaign Writer**: Write effective marketing emails from subject line to CTA.
    - **Blog Generator**: Generate well-structured and SEO-friendly blog posts.
    - **SEO Analyst**: Get keyword ideas, competitor analysis, and on-page SEO tips.
    - **Content Improver**: Rewrite your existing text to be more persuasive, clear, or professional.
    - **AI to Human Text Converter**: Make AI-generated text sound more natural and human-like.
    - **Digital Marketing Analyst**: Analyze data and get actionable insights from your reports.
    """)

# --- Sidebar ---
with st.sidebar:
    st.header("Tools & Settings")
    
    # Tool selection
    selected_tool = st.selectbox("Choose your marketing tool:", list(SYSTEM_PROMPTS.keys()))
    
    st.markdown("---")
    
    # File Uploader
    st.subheader("Upload a File")
    uploaded_file = st.file_uploader(
        "Upload a file for context (not used for Image Generation).",
        type=['png', 'jpg', 'jpeg', 'docx', 'pptx', 'xlsx', 'sql', 'wav', 'mp3', 'ogg']
    )

    st.markdown("---")

    # Clear Chat Button with Confirmation
    if 'confirm_clear' not in st.session_state:
        st.session_state.confirm_clear = False

    if st.button("Clear Chat History"):
        st.session_state.confirm_clear = True

    if st.session_state.confirm_clear:
        st.warning("Are you sure you want to clear the chat history?")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Yes, Clear It"):
                st.session_state.messages = []
                st.session_state.confirm_clear = False
                st.rerun()
        with col2:
            if st.button("Cancel"):
                st.session_state.confirm_clear = False
                st.rerun()


# --- Chat History Initialization ---
if "messages" not in st.session_state:
    st.session_state.messages = []

# Display chat messages from history
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        if "image" in message:
            st.image(message["image"], width=200)
        if "generated_image" in message:
            st.image(message["generated_image"], caption="Generated Image")
        if "content" in message:
            st.markdown(message["content"])

# --- Chat Input and Logic ---
if prompt := st.chat_input("What can I help you with today?"):
    user_message = {"role": "user", "content": prompt}
    
    if uploaded_file is not None and selected_tool != "Image Generator":
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        
        if file_extension in ['.png', '.jpg', '.jpeg']:
            image = Image.open(uploaded_file)
            user_message["image"] = image
        # Other file handling logic...
        
    st.session_state.messages.append(user_message)
    with st.chat_message("user"):
        if "image" in user_message:
            st.image(user_message["image"], width=200)
        st.markdown(user_message["content"])

    # --- Generate AI Response ---
    with st.chat_message("assistant"):
        with st.spinner("🤖 Thinking..."):
            if selected_tool == "Image Generator":
                generated_image = generate_image(prompt)
                if isinstance(generated_image, Image.Image):
                    st.image(generated_image, caption="Generated Image")
                    st.session_state.messages.append({"role": "assistant", "generated_image": generated_image})
                else:
                    st.error(generated_image) # Display error message
                    st.session_state.messages.append({"role": "assistant", "content": generated_image})
            else:
                # Text generation logic
                full_prompt = [SYSTEM_PROMPTS[selected_tool]]
                for msg in st.session_state.messages:
                    if "content" in msg:
                        full_prompt.append(f"{msg['role']}: {msg['content']}")

                model_input = []
                if "image" in user_message:
                    model_input.append(user_message["image"])
                model_input.append("\n".join(full_prompt))

                try:
                    response = text_model.generate_content(model_input)
                    response_text = response.text
                except Exception as e:
                    response_text = f"Sorry, an error occurred: {e}"
                
                st.markdown(response_text)
                st.session_state.messages.append({"role": "assistant", "content": response_text})
